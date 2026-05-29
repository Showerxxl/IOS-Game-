"""
Переделывает слайд 7 и добавляет 2 слайда (интеграция DQN и LSTM) по шаблону слайда 7.
Сохраняет в новый файл, оригинал не трогает.
"""
import copy
import os
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor

SRC = "/Users/mac/Desktop/term/presentarion.pptx"
OUT = "/Users/mac/Desktop/term/presentarion_ml.pptx"
ART = "/Users/mac/Desktop/term/Котозрыв/ML/artifacts"

TITLE_COLOR = RGBColor(0x10, 0x2D, 0x69)


def find_shape(slide, name):
    for sh in slide.shapes:
        if sh.name == name:
            return sh
    return None


def set_title(slide, text):
    sh = find_shape(slide, "Текст 4")
    if sh and sh.has_text_frame and sh.text_frame.paragraphs[0].runs:
        sh.text_frame.paragraphs[0].runs[0].text = text
    return sh


def set_body(shape, paragraphs, size=14):
    """paragraphs: список строк-абзацев."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, line in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)


R_NS = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"


def _remap_rids(spTree, rid_map):
    for el in spTree.iter():
        for attr in list(el.attrib):
            if attr.startswith(R_NS) and el.attrib[attr] in rid_map:
                el.attrib[attr] = rid_map[el.attrib[attr]]


def duplicate_slide(prs, index):
    source = prs.slides[index]
    layout = source.slide_layout
    new_slide = prs.slides.add_slide(layout)
    # убрать плейсхолдеры, которые добавил layout
    for shp in list(new_slide.shapes):
        shp._element.getparent().remove(shp._element)
    # скопировать все фигуры исходного слайда
    for shp in source.shapes:
        new_slide.shapes._spTree.append(copy.deepcopy(shp._element))
    # перенести связи (картинки и пр.), пропуская layout/notes; rId переотображаем
    rid_map = {}
    for rId, rel in source.part.rels.items():
        if "slideLayout" in rel.reltype or "notesSlide" in rel.reltype:
            continue
        new_rId = new_slide.part.relate_to(rel._target, rel.reltype, rel.is_external)
        rid_map[rId] = new_rId
    _remap_rids(new_slide.shapes._spTree, rid_map)
    return new_slide


def move_slide_after(prs, src_index, after_index):
    """Перемещает слайд src_index так, чтобы он шёл сразу после after_index."""
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    el = ids[src_index]
    sldIdLst.remove(el)
    # после удаления индексы сдвинулись; вставляем по позиции after_index+1
    sldIdLst.insert(after_index + 1, el)


def remove_pictures(slide):
    for sh in list(slide.shapes):
        if sh.shape_type == 13:  # PICTURE
            sh._element.getparent().remove(sh._element)


IN = 914400  # EMU в дюйме


def set_pos(shape, left=None, top=None, width=None, height=None):
    if left is not None:   shape.left = Emu(int(left * IN))
    if top is not None:    shape.top = Emu(int(top * IN))
    if width is not None:  shape.width = Emu(int(width * IN))
    if height is not None: shape.height = Emu(int(height * IN))


def clear_dot(slide):
    sh = find_shape(slide, "Текст 1")
    if sh and sh.has_text_frame:
        sh.text_frame.clear()


def add_centered_picture(slide, path, top, max_h):
    """Добавляет картинку по центру по горизонтали, вписывая высоту в max_h."""
    pic = slide.shapes.add_picture(path, Emu(0), Emu(int(top * IN)))
    ratio = pic.height / pic.width
    h = max_h
    w = h / ratio
    pic.height = Emu(int(h * IN))
    pic.width = Emu(int(w * IN))
    pic.left = Emu(int(((13.3 - w) / 2) * IN))
    pic.top = Emu(int(top * IN))
    return pic


def main():
    prs = Presentation(SRC)

    # --- 1) Переделать слайд 7 (index 6): обзор ML-подхода (текстовый слайд) ---
    s7 = prs.slides[6]
    set_title(s7, "Алгоритм работы ИИ: машинное обучение")
    clear_dot(s7)
    remove_pictures(s7)  # убираем скриншоты старой эвристики
    body7 = find_shape(s7, "TextBox 12")
    set_pos(body7, left=0.58, top=2.45, width=12.1, height=4.2)
    set_body(body7, [
        "Рукописная эвристика заменена двумя моделями машинного обучения, "
        "работающими на устройстве через Core ML.",
        "• DQN-агент (обучение с подкреплением) — играет за противника.",
        "• LSTM-сеть — «личный коуч», разбирает партию игрока после её завершения.",
        "Единый пайплайн: собственный симулятор игры на Python → обучение в PyTorch → "
        "конвертация в Core ML → инференс на устройстве.",
        "Кодировщик состояния на Swift точно повторяет Python — это гарантирует "
        "совпадение обучения и рантайма.",
    ], size=18)

    # --- 2) Дублировать слайд 7 дважды (новые слайды уходят в конец) ---
    dqn = duplicate_slide(prs, 6)
    lstm = duplicate_slide(prs, 6)

    # --- 3) Заполнить слайд DQN ---
    set_title(dqn, "Интеграция ML-бота (DQN)")
    clear_dot(dqn)
    body_dqn = find_shape(dqn, "TextBox 12")
    set_pos(body_dqn, left=0.58, top=1.45, width=12.1, height=1.9)
    set_body(body_dqn, [
        "Противник — нейросеть Deep Q-Network: вход — 31 признак состояния (рука, "
        "опасность колоды, знание верхних карт), выход — Q-оценки 8 действий.",
        "Обучение: self-play, 400 000 шагов; 3 чекпойнта = 3 уровня сложности.",
        "На устройстве MLBotService кодирует состояние, маскирует нелегальные ходы и "
        "берёт argmax. Итог: ML-бот обыгрывает прежнюю эвристику в ~68% партий.",
    ], size=13)
    remove_pictures(dqn)
    add_centered_picture(dqn, os.path.join(ART, "winrate_curve.png"), top=3.45, max_h=3.7)

    # --- 4) Заполнить слайд LSTM ---
    set_title(lstm, "Интеграция коуча (LSTM)")
    clear_dot(lstm)
    body_lstm = find_shape(lstm, "TextBox 12")
    set_pos(body_lstm, left=0.58, top=1.45, width=12.1, height=1.9)
    set_body(body_lstm, [
        "После партии LSTM строит «кривую вероятности победы» по ходам игрока "
        "(вход — последовательность T × 39).",
        "Данные размечены DQN-оракулом: 40 000 партий; accuracy 94%, ROC-AUC 0.99.",
        "CoachService находит ходы с резким падением вероятности и сравнивает с советом "
        "DQN: формирует подсказки и определяет стиль игры. Разбор — на экране CoachAnalysisView.",
    ], size=13)
    remove_pictures(lstm)
    add_centered_picture(lstm, os.path.join(ART, "coach_winprob_example.png"), top=3.45, max_h=3.7)

    # --- 5) Переставить новые слайды сразу после слайда 7 ---
    # сейчас порядок: ... s7(6), ... , dqn(предпоследний), lstm(последний)
    n = len(prs.slides._sldIdLst)
    move_slide_after(prs, n - 2, 6)   # dqn -> позиция 7 (index 7)
    move_slide_after(prs, n - 1, 7)   # lstm -> позиция 8 (index 8)

    prs.save(OUT)
    print("Сохранено:", OUT)
    print("Всего слайдов:", len(prs.slides._sldIdLst))


if __name__ == "__main__":
    main()
